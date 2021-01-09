from pptx import Presentation
#For Presentation
from pptx.util import Pt
#For text formatting
import  numpy as np, pandas as pd
#For excel
import os, copy
from time import perf_counter

t1 = perf_counter()
prs_input_path = r"C:\Users\adam.smith04\Documents\PROJECTS\mu_code\wavemark\resources\QBR template 2021.pptx"
df_file_path = r"C:\Users\adam.smith04\Documents\PROJECTS\mu_code\wavemark\resources\QBR Template FY21 v0.1.xlsx"

#Create data frame
df = pd.read_excel(df_file_path)

#creates a unique list of PAC RDs
idn_list = list(df['IDN_WaveMark'].unique())

#set presentation you want to work with
prs = Presentation(prs_input_path)

#Add slides to a list
slides = list(prs.slides)

#give me the first customer in the list
for customer in idn_list:
    #start with the 5th slide in the deck
    for slide in slides[4:15]:
        #give me the first shape in the group of shapes on the slide
        for shape in slide.shapes:
            #if the shape is the textbox we want, add the text based on slide slide number and format all text the same
            if shape.has_text_frame and shape.name == 'TextBox 6':
                if slides.index(slide) in [4, 5, 6, 7]:
                    shape.text = f'Total inventory, {customer} enterprise'
                
                elif slides.index(slide) in [8, 9]:
                    shape.text = f'Expiration, {customer} enterprise'

                elif slides.index(slide) in [10, 11]:
                    shape.text = f'Wasted supplies, {customer} enterprise'
                
                elif slides.index(slide) in [12, 13]:
                    shape.text = f'Accurate supply documentation, {customer} enterprise'

                #to get the formatting of the text you have to 
                #add the text to the run through the text_frame
                text_frame = shape.text_frame
                p = text_frame.paragraphs[0]
                p.font.name = 'Arial (Body)'
                p.font.size = Pt(12)
                p.font.italic = True

    print(f'Saving {customer} QBR powerpoint file...')
    prs.save(f'C:/Users/adam.smith04/Documents/PROJECTS/mu_code/wavemark/outputs/QBR template {customer} 2021.pptx')       

t2 = perf_counter()
print('Took {} Seconds'.format(t2-t1))
